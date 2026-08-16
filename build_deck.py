#!/usr/bin/env python3
"""Agent-native accounting deck.

Vantage: we are a tech startup that has been automating existing accounting software
with RPA, now rebuilding it agent-native. Audience: technically-literate investors.

Versioning is handled by git — this script writes one canonical deck file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- design tokens

INK    = RGBColor(0x11, 0x1A, 0x27)
BODY   = RGBColor(0x33, 0x3E, 0x4F)
MUTED  = RGBColor(0x6B, 0x76, 0x88)
ACCENT = RGBColor(0x0E, 0x74, 0x90)
WARM   = RGBColor(0xB4, 0x53, 0x09)
DEEPW  = RGBColor(0x7A, 0x3B, 0x08)
DEEPA  = RGBColor(0x0A, 0x53, 0x66)
PANEL  = RGBColor(0xF3, 0xF6, 0xF8)
PANEL2 = RGBColor(0xE8, 0xEE, 0xF2)
LINE   = RGBColor(0xD5, 0xDC, 0xE4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TINT_A = RGBColor(0xE4, 0xF0, 0xF4)
TINT_W = RGBColor(0xFA, 0xEE, 0xE3)

FONT = "Helvetica Neue"
SW, SH = 13.333, 7.5
ML, MR = 0.72, 0.72
CW = SW - ML - MR

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
_page = {"n": 0}


# ------------------------------------------------------------------- primitives

def new_slide(numbered=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if numbered:
        _page["n"] += 1
        tb = s.shapes.add_textbox(Inches(SW - MR - 0.8), Inches(SH - 0.36),
                                  Inches(0.8), Inches(0.26))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(_page["n"])
        r.font.size, r.font.name, r.font.color.rgb = Pt(10), FONT, LINE
    return s


def text(slide, s, x, y, w, h, size=15, color=BODY, bold=False, align=PP_ALIGN.LEFT,
         space=6, line=1.25, caps=False, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space)
    p.line_spacing = line
    r = p.add_run()
    r.text = s.upper() if caps else s
    r.font.size, r.font.name, r.font.bold = Pt(size), FONT, bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def eyebrow(slide, label):
    text(slide, label, ML, 0.46, CW, 0.28, size=10.5, color=ACCENT, bold=True, caps=True)


def heading(slide, s, size=27, y=0.82, w=None):
    return text(slide, s, ML, y, w or CW, 1.0, size=size, color=INK, bold=True, line=1.1)


def box(slide, x, y, w, h, fill=PANEL, edge=None, radius=0.04):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if edge:
        sh.line.color.rgb = edge; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.text_frame.text = ""
    return sh


def bullets(slide, items, x, y, w, h, size=15, gap=9, line=1.28):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for body, kind in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line
        if kind == "h":
            p.space_before = Pt(gap + 5); p.space_after = Pt(2)
            r = p.add_run(); r.text = body
            r.font.size, r.font.bold, r.font.color.rgb = Pt(size), True, INK
        elif kind == "b":
            p.space_after = Pt(gap)
            r = p.add_run(); r.text = "— " + body
            r.font.size, r.font.color.rgb = Pt(size), BODY
        elif kind == "plain":
            p.space_after = Pt(gap)
            r = p.add_run(); r.text = body
            r.font.size, r.font.color.rgb = Pt(size), BODY
        else:
            p.space_before = Pt(gap); p.space_after = Pt(2)
            r = p.add_run(); r.text = body
            r.font.size, r.font.italic, r.font.color.rgb = Pt(size - 2), True, WARM
        r.font.name = FONT
    return tb


def table(slide, data, x, y, w, col_w, row_h=0.42, head_h=0.44, size=11.5,
          emphasize_col=None):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w),
                               Inches(head_h + row_h * (rows - 1))).table
    gt.first_row = False
    gt.horz_banding = False
    for i, cwi in enumerate(col_w):
        gt.columns[i].width = Inches(cwi)
    gt.rows[0].height = Inches(head_h)
    for i in range(1, rows):
        gt.rows[i].height = Inches(row_h)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = INK
            elif emphasize_col is not None and ci == emphasize_col:
                cell.fill.fore_color.rgb = TINT_W
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else PANEL
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            r = p.add_run(); r.text = val
            r.font.size, r.font.name = Pt(size), FONT
            if ri == 0:
                r.font.bold = True; r.font.color.rgb = WHITE
            elif ci == 0:
                r.font.bold = True; r.font.color.rgb = INK
            elif emphasize_col is not None and ci == emphasize_col:
                r.font.color.rgb = DEEPW
            else:
                r.font.color.rgb = BODY
    return gt


def notes(slide, body):
    slide.notes_slide.notes_text_frame.text = body.strip()


def card(slide, x, y, w, h, title_, body, fill=PANEL, edge=LINE, tcol=INK,
         tsize=12.5, bsize=10.5):
    box(slide, x, y, w, h, fill=fill, edge=edge)
    text(slide, title_, x + 0.24, y + 0.14, w - 0.46, 0.28, size=tsize, color=tcol, bold=True)
    text(slide, body, x + 0.24, y + 0.45, w - 0.46, h - 0.58, size=bsize, color=BODY, line=1.32)


# ============================================================== 1. TITLE

s = new_slide(numbered=False)
box(s, 0, 0, SW, SH, fill=INK, radius=0.0)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Pt(6))
band.fill.solid(); band.fill.fore_color.rgb = WARM
band.line.fill.background(); band.shadow.inherit = False

text(s, "Aug 2026", ML, 1.55, CW, 0.3,
     size=11.5, color=RGBColor(0x7A, 0xB8, 0xC8), bold=True, caps=True)
text(s, "Agent-Native\nAccounting Software", ML, 2.15, 11.4, 2.4, size=46, color=WHITE,
     bold=True, line=1.08)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(5.2), Inches(2.2), Pt(2.5))
ln.fill.solid(); ln.fill.fore_color.rgb = WARM
ln.line.fill.background(); ln.shadow.inherit = False
text(s, "Chunliang Lyu", ML, 5.5, CW, 0.4, size=14, color=RGBColor(0x8E, 0x9C, 0xAE))

notes(s, """
Opening line: "There are two ways to do AI in accounting. One is to bolt a copilot onto the
software that already exists. The other is to rebuild the software so it produces the work
product instead of recording it. Only the second one changes anything, and it needs a
different technical foundation — which is what this deck is about."

Say where we come from early, because it is the credibility. We already automate this
software for a living. Our RPA bots drive the packages our customers use, which means we
have seen exactly which steps are mechanical, exactly where a human has to intervene, and
exactly how brittle the seams are. Nobody had to tell us where the work is.

Structure: technical background (2 slides), why this domain is uniquely suited (1), the
architecture and its four load-bearing components (6), the two hard problems (1), then how
we build and the implementation plan. Tell them the architecture section is the real content
and to interrupt there.
""")


# ============================================================== 2. THESIS

s = new_slide()
eyebrow(s, "The thesis")
heading(s, "Software stopped at the edge of the work.\nThat edge was a technical limit, and it moved.", size=25)

w3 = 3.79
xs = ML
thesis = [
    ("The constraint", PANEL, INK,
     "Deterministic code cannot act under ambiguity. Every path has to be enumerated in "
     "advance, and there is no path for \"it depends\". So every judgment was routed out of "
     "the software and into a person."),
    ("What lifted", PANEL, INK,
     "Models tolerate ambiguity, cover the long tail at near-zero marginal cost, and need no "
     "interface built in advance. Three specific engineering properties — not a claim that "
     "models are clever."),
    ("Why us", TINT_W, DEEPW,
     "We already automate this software with RPA. Our customers never cared which package "
     "sat underneath — they wanted the result. RPA automated the clicking; it never "
     "automated the deciding. We know exactly where that line falls."),
]
for name, fill, tcol, body in thesis:
    box(s, xs, 2.35, w3, 3.15, fill=fill, edge=WARM if fill is TINT_W else LINE)
    text(s, name, xs + 0.28, 2.6, w3 - 0.55, 0.4, size=17.5, color=tcol, bold=True)
    text(s, body, xs + 0.28, 3.15, w3 - 0.55, 2.1, size=13, color=BODY, line=1.42)
    xs += w3 + 0.28

box(s, ML, 5.75, CW, 1.1, fill=INK)
text(s, "The bet is not that models are smart. It is that accounting is verifiable — so we "
        "can prove our output is correct before it ever touches a client's books.",
     ML + 0.34, 6.05, CW - 0.68, 0.6, size=16, color=WHITE, bold=True, line=1.3)

notes(s, """
This is the whole argument in one slide. Deliver it, pause, then say the rest of the deck
defends the third box and the dark bar.

The first box is the point most people miss: the reason accounting software makes you fill in
forms is not laziness or bad design. Deterministic code physically cannot handle "it
depends", so the ambiguity had to go somewhere, and it went to the user. Every hour of
bookkeeper time is the cost of that architectural constraint.

The third box is our origin story and it is the best "why us" we have. We have spent years
driving these packages with RPA, which taught us two things nobody else knows as precisely.
First, exactly where the mechanical work stops and judgment starts — our bots hit that wall
thousands of times a day. Second, that customers are completely indifferent to the software
underneath; they asked us for an outcome and treated the package as plumbing. That second
point is the permission to replace it. If customers were attached to the software, this
would be a much harder pitch.

Be honest about RPA's ceiling if asked: it automates the keystrokes around a fixed UI, it
breaks when the UI changes, and it cannot make a judgment call. It was the best available
answer under the old constraint. We are not abandoning it out of fashion — the constraint
lifted.

The dark bar is the sentence to repeat if you get only one. Almost every agent pitch in the
market cannot say it, because almost no domain lets you check the answer mechanically. Ours
does.

If someone jumps ahead to "why won't the incumbents just ship this" — park it, you answer it
structurally on the architecture slides. The short version: the substrate they would need is
a schema decision they cannot retrofit, and their revenue depends on the seats this removes.
""")


# ============================================================== 3. WHY THE SOFTWARE IS SHAPED THIS WAY

s = new_slide()
eyebrow(s, "Technical background · 1 of 2")
heading(s, "Why our accounting software is shaped the way it is")

text(s, "Not a criticism of the vendors — a consequence of what deterministic software can "
        "and cannot do. Four design choices follow from one limit, and all four put the work "
        "on us.",
     ML, 1.9, 11.2, 0.6, size=15, color=BODY, line=1.35)

cols = [
    ("Pre-specified paths",
     "Every branch must be enumerated by a developer before it runs. There is no code path "
     "for \"this invoice is ambiguous\". Anything requiring interpretation has to leave the "
     "program."),
    ("The form as adapter",
     "A form converts fuzzy human intent into a structure the program can execute. That is "
     "its actual function. Forms are the scar tissue of the constraint above — which is why "
     "there are so many of them."),
    ("Validation is syntactic",
     "The software checks that the date is a date and the entry balances. It does not check "
     "whether the expense belongs in that account. Well-formed and correct are different "
     "questions, and it only asks the first."),
    ("The human is the error handler",
     "Semantic correctness was outsourced to us. The vendor's contract is: prevent invalid "
     "states, show a clear error. A far cheaper contract than being right — and the reason "
     "our cost structure is people."),
]
cx = ML
for h, b in cols:
    w = 2.79
    text(s, h, cx, 2.72, w, 0.62, size=16, color=ACCENT, bold=True, line=1.15)
    text(s, b, cx, 3.48, w, 2.3, size=12.5, color=BODY, line=1.42)
    cx += w + 0.25

box(s, ML, 5.95, CW, 1.0, fill=TINT_W)
text(s, "Our software was never trying to do the accounting. It was trying to stop us from "
        "entering something impossible. Everything in this deck follows from taking on the "
        "harder contract: being right, and proving it.",
     ML + 0.32, 6.2, CW - 0.64, 0.6, size=14.5, color=DEEPW, bold=True, line=1.3)

notes(s, """
This slide replaces the business-model background from the earlier draft. It does the same
setup job but in technical terms, which suits both the audience and the fact that we are
insiders describing tools we use daily.

The third column is the one to slow down on, because it is the pivot for the entire
architecture section. Say it plainly: "Every accounting system we have ever used validates
syntax, not semantics. It will happily let us post a perfectly balanced entry to completely
the wrong account. Well-formed and correct are different questions, and no software we own
asks the second one." Then: the whole point of what we are building is a system that asks
the second question, and can answer it.

The fourth column is where the audience feels it. Everyone in the room who has staffed a
close knows that the software's job ends exactly where the expensive part begins.

Watch the tone — do not bash Xero or QBO. Their design was correct for the constraint they
faced. Being fair about that makes the claim that the constraint has lifted more credible,
not less.
""")


# ============================================================== 4. THREE CONSTRAINTS

s = new_slide()
eyebrow(s, "Technical background · 2 of 2")
heading(s, "Three constraints forced that design. All three just lifted.")

hdr = [["The constraint", "What it forced", "What changes now"]]
rows = [
    ["No tolerance for ambiguity",
     "Every path pre-specified. Anything needing interpretation of unstructured input was "
     "handed to a person behind a form.",
     "Models absorb messy documents, bank memos, contracts and email threads — the exact "
     "inputs the software pushed onto a human operator."],
    ["No economics for the long tail",
     "The 1,000th edge case costs as much to build as the first. Vendors served the common "
     "case and left the tail to us and to spreadsheets.",
     "Marginal cost of covering an edge case approaches the cost of describing it. The tail "
     "becomes addressable for the first time."],
    ["Every capability needed a pre-built interface",
     "We could only direct the software through affordances a product manager imagined and "
     "a roadmap funded.",
     "Natural language is a universal interface. Nobody has to anticipate the request in "
     "advance for it to be answerable."],
]
table(s, hdr + rows, ML, 2.05, CW, [3.3, 4.29, 4.3], row_h=1.20, head_h=0.46,
      size=12.5, emphasize_col=2)

box(s, ML, 6.18, CW, 0.85, fill=TINT_A)
text(s, "This is the technical claim under the whole strategy — three specific properties, "
        "not a general belief about AI: ambiguity tolerance, near-zero marginal cost on the "
        "long tail, and an interface nobody has to build in advance.",
     ML + 0.32, 6.38, CW - 0.64, 0.5, size=14, color=DEEPA, bold=True, line=1.3)

notes(s, """
This slide earns credibility precisely by refusing the hype framing. Any investor who has
sat through forty AI pitches has heard "AI changes everything"; very few have heard a
precise statement of which constraints lifted and which did not.

Walk the rows left to right, one sentence each. Then deliver the bottom box close to
verbatim — it is the most quotable line in the deck.

Expected pushback: "the long-tail claim is doing a lot of work." Concede it partially and
immediately — marginal cost on the tail falls a great deal but not to zero, and reliability
on tail cases is measurably worse than on the common case. That is exactly why the
architecture has a verification layer and a human review queue rather than assumed
autonomy. Point forward to the verification slide; do not defend more than the evidence
supports.
""")


# ============================================================== 5. THE ORACLE

s = new_slide()
eyebrow(s, "Domain fit")
heading(s, "Accounting is the best domain for this, and the reason is verification")

text(s, "Most agent products cap out at \"assistive\" because nobody can mechanically tell "
        "whether the output is right. In law, in consulting, in most knowledge work there is "
        "no oracle. Accounting hands us one.",
     ML, 1.88, 11.4, 0.6, size=15, color=BODY, line=1.35)

w3 = 3.79
xs = ML
panels = [
    ("1 · The domain gives us a verifier", TINT_A, DEEPA,
     ["Debits must equal credits on every entry",
      "Trial balance must tie",
      "Subledgers must tie to control accounts",
      "Bank reconciliation must net to zero",
      "Tax and filing rules are codified, not inferred"], None),
    ("2 · So we can verify before we commit", TINT_W, DEEPW, None,
     "That single capability is the difference between a demo and a system we would put in "
     "front of a client.\n\nIt lets us auto-post with a measured error bound, route only "
     "genuine ambiguity to a person, and state our own accuracy in terms an external "
     "reviewer will accept."),
    ("3 · And we know where the work is", TINT_W, DEEPW, None,
     "Years of driving these packages with RPA is a map of exactly which steps are mechanical "
     "and where a human has to intervene.\n\nOur bots hit that wall thousands of times a day, "
     "and the logs record it. Combined with customer closes we can access through existing "
     "engagements, that is a ground-truth corpus with a running start."),
]
for name, fill, tcol, blist, bodytext in panels:
    box(s, xs, 2.62, w3, 3.12, fill=fill, edge=WARM if fill is TINT_W else ACCENT)
    text(s, name, xs + 0.26, 2.85, w3 - 0.5, 0.32, size=13.5, color=tcol, bold=True)
    if blist:
        bullets(s, [(b, "b") for b in blist], xs + 0.26, 3.28, w3 - 0.5, 2.2, size=12.5, gap=7)
    else:
        text(s, bodytext, xs + 0.26, 3.28, w3 - 0.5, 2.3, size=12.5, color=BODY, line=1.42)
    xs += w3 + 0.28

box(s, ML, 5.98, CW, 0.95, fill=INK)
text(s, "An oracle plus a map of where the judgment actually sits is not a nice-to-have. It "
        "is the reason autonomy is shippable in accounting before almost anywhere else — and "
        "the reason it is shippable by us rather than by someone starting cold.",
     ML + 0.34, 6.2, CW - 0.68, 0.55, size=15, color=WHITE, bold=True, line=1.3)

notes(s, """
This is the slide the audience should leave remembering, and it answers two questions at
once: why this vertical, and why us rather than a well-funded generalist.

Make the general point first. In most knowledge work you cannot check the answer
mechanically, so a human never leaves the loop and the product ceiling is assistive. We did
not pick accounting because it is exciting. We picked it because it is verifiable.

Then panel 3, which is our unfair advantage. Spell out the RPA angle: for years our bots
have executed the mechanical steps and stopped dead at the judgment ones. Every one of those
stops is a logged, timestamped record of exactly where deterministic automation runs out.
A team starting cold spends a year discovering that map; we have been maintaining it as a
production system. It also tells us which categories to attack first, in what order, by
volume and by breakage rate.

Be precise about what we do and do not have. We have process telemetry and exception logs.
We do not own our customers' books — access for training and evaluation comes through
existing engagements, with consent, and some agreements will need updating. Volunteer that
before it is asked; it is a process question, not a blocker, and pretending otherwise is
worse than the caveat.

If asked to quantify: have real numbers ready — customers, bots in production, transactions
touched per month, intervention rate. That is the most concrete asset claim we can make.
""")


# ============================================================== 6. WHAT BECOMES DEFENSIBLE

s = new_slide()
eyebrow(s, "What this changes")
heading(s, "Where the defensible engineering sits now")

text(s, "When the vendor's differentiation was the interface, the moat was habit and "
        "configuration. When the system does the work, four different things become hard to "
        "copy — and none of them is the prompt.",
     ML, 1.88, 11.4, 0.6, size=15, color=BODY, line=1.35)

quads = [
    ("Verification", "Anyone can prompt a model. Almost nobody can prove the output is "
     "correct. In a domain with an oracle, the verifier is the product."),
    ("Context", "Accumulated client-specific policy and history: this vendor is always COGS, "
     "this client capitalises above [$2,500]. Proprietary state, not a model."),
    ("Permission", "Write access to the ledger, the bank, the filing. Reversing that is a "
     "governance decision with a liability trail — a far stronger lock-in than UI habit."),
    ("Evals", "A corpus of verified outcomes to regression-test against. Prompts copy in an "
     "afternoon; a decade of reviewed closes does not."),
]
xq, yq = ML, 2.62
for i, (h, b) in enumerate(quads):
    x = ML + (i % 2) * 3.02
    y = 2.62 + (i // 2) * 1.62
    text(s, h, x, y, 2.78, 0.3, size=16, color=ACCENT, bold=True)
    text(s, b, x, y + 0.4, 2.78, 1.1, size=12.5, color=BODY, line=1.4)

box(s, ML + 6.35, 2.62, 5.55, 3.12, fill=PANEL)
text(s, "What does not change", ML + 6.65, 2.86, 5.0, 0.3, size=11, color=MUTED,
     bold=True, caps=True)
text(s, "The system of record matters more, not less. Agents need a substrate with truth, "
        "permissions and history — they make interfaces obsolete, not databases. Anyone who "
        "tells you the ledger goes away has not built one.\n\nTrust, security, "
        "confidentiality and professional standards still gate everything. Our customers' books "
        "our reputation are the assets at risk, which is why the controls in this "
        "architecture are not optional extras.",
     ML + 6.65, 3.24, 5.0, 2.4, size=12.5, color=BODY, line=1.44)

box(s, ML, 5.98, CW, 0.9, fill=TINT_W)
text(s, "Note what is not on this list: the agent layer. Prompts are the cheapest and most "
        "replaceable component in the system, and we have deliberately built the least of it.",
     ML + 0.32, 6.2, CW - 0.64, 0.5, size=14, color=DEEPW, bold=True, line=1.3)

notes(s, """
This is the bridge into the architecture. It tells the audience what to look for in the next
six slides.

Draw out permission specifically, because it is the one an investor will underrate. UI habit
is a preference and preferences get overridden in a procurement cycle. Write access to a
client's ledger and bank, with our professional liability attached, is a governance decision
with a documented approval trail. Reversing it requires a partner-level conversation, not a
preference change.

The "what does not change" panel is deliberate and it buys more credibility than anything
else on the slide. Everyone in this category has heard someone claim agents make databases
or systems of record obsolete. Saying plainly that the ledger becomes MORE important signals
we have actually built something.

The bottom bar sets up the biggest surprise in the architecture section — that the agent
layer is one of the smallest line items. Plant it here so it lands twice.

If asked which of the four we have today: be honest. Verification is built, evals are in
progress, context accrues per client per month, permission comes with each engagement. Do
not claim all four are mature.
""")


# ============================================================== 7. ARCHITECTURE

s = new_slide()
eyebrow(s, "Architecture")
heading(s, "The agent decides. Deterministic code executes.\nNothing commits unverified.", size=24)

LX, LW = ML, 9.85
XX, XW = ML + 10.05, 1.83
top = 2.42
bh, bg = 0.515, 0.075

layers = [
    ("Client surfaces", "exception queue (the real UI) · chat & email · client reporting", PANEL, False),
    ("Verification & guardrails", "hard invariants · anomaly and variance · confidence routing", TINT_W, True),
    ("Agent layer", "workflow supervisors · extract · classify · match · chase · analyse", PANEL, False),
    ("Durable orchestration", "the close checklist as a resumable state machine, not free-form agency", PANEL, False),
    ("Deterministic tool layer", "typed, invariant-enforcing actions: post_entry · reconcile · close_period", TINT_W, True),
    ("Ledger core", "append-only · bitemporal · double-entry · full provenance", TINT_W, True),
    ("Ingestion", "banks · cards · payroll · POS · AP inbox · documents · prior-year books", PANEL, False),
]
y = top
for name, detail, fill, star in layers:
    box(s, LX, y, LW, bh, fill=fill, edge=WARM if fill is TINT_W else LINE)
    text(s, name, LX + 0.22, y + 0.075, 3.0, 0.36, size=13.5, color=INK, bold=True)
    text(s, detail, LX + 3.3, y + 0.095, LW - 3.55, 0.36, size=11.5, color=BODY)
    if star:
        text(s, "◆", LX + LW - 0.34, y + 0.09, 0.3, 0.3, size=12, color=WARM, bold=True)
    y += bh + bg

box(s, XX, top, XW, y - top - bg, fill=PANEL2, edge=LINE)
text(s, "Cross-cutting", XX + 0.16, top + 0.14, XW - 0.32, 0.3, size=10, color=MUTED,
     bold=True, caps=True)
for label, dy in [("Audit trail\n& lineage", 0.6), ("Evals &\nreplay harness", 1.3),
                  ("Cost & token\ntelemetry", 2.0), ("Agent identity,\nleast privilege", 2.7)]:
    text(s, label, XX + 0.16, top + dy, XW - 0.32, 0.6, size=11, color=BODY, line=1.25)

text(s, "◆  Proprietary and defensible. The agent layer is the part we will rewrite most "
        "often and depend on least.",
     ML, y + 0.1, CW, 0.4, size=12.5, color=WARM, bold=True)

notes(s, """
Spend the most time here. Read the title as a sentence — it is the design philosophy in nine
words, and it is what makes the output defensible to a customer's auditor.

Walk it bottom-up, not top-down: ingestion, ledger, tools, orchestration, agents,
verification, surfaces. Bottom-up is the build order and it makes the dependency structure
obvious.

The point to hammer: the three diamonds are where the value is, and the agent layer is not
one of them. Teams in this category over-invest in the agent layer because it is the
visible, demo-able part, then find they have no verification and no evals and cannot tell
whether they are improving. We inverted that deliberately.

Second point, and it is the one a technical investor will test: verification sits ABOVE the
agent and BELOW the client. Nothing the agent produces reaches a client's books without
passing deterministic checks first. That is an architectural guarantee, not a policy or a
promise about model behaviour.

Expect: "why durable orchestration rather than letting the agent plan the close?" Answer: a
close is a known checklist that spans days, waits on client email, and needs compensation
logic when a step fails. We model the known process deterministically and use agents only
inside the judgment-bearing steps. Free-roaming agency over a general ledger is how you get
a restatement you cannot explain.
""")


# ============================================================== 8. LEDGER

s = new_slide()
eyebrow(s, "Key component · 1 of 4")
heading(s, "Ledger substrate: append-only, bitemporal, fully attributed")

text(s, "Ordinary software updates a row and moves on. Three requirements that no accounting "
        "package we use has make that impossible for us — and all three are schema decisions, "
        "so all three must be made on day one. They cannot be retrofitted.",
     ML, 1.9, 11.3, 0.7, size=14.5, color=BODY, line=1.35)

items = [
    ("Bitemporality — booking date versus effective date", "h"),
    ("We must be able to answer \"what did the books say as of 31 March?\" for review and "
     "for audit. Critically, it is also what lets us replay an agent against a historical "
     "period and score it. Our entire eval capability depends on this one schema choice.", "b"),
    ("Immutability — corrections are reversing entries, never updates", "h"),
    ("A mutable ledger cannot be explained after the fact, and an agent that can silently "
     "overwrite is not defensible at any price.", "b"),
    ("Provenance as a first-class, queryable chain", "h"),
    ("source document → extraction → reasoning trace → tool call → policy applied → "
     "approver → posting. In our current stack an audit log is a compliance artefact bolted "
     "on the side. Here it is load-bearing: it is how the work product gets defended.", "b"),
    ("Every actor is typed — human, agent or system — with the agent's identity, model "
     "version and prompt revision recorded on the entry itself.", "note"),
]
bullets(s, items, ML, 2.78, 7.4, 3.9, size=13.5)

box(s, ML + 7.85, 2.78, 4.02, 3.5, fill=PANEL)
text(s, "Why this is the slide that matters", ML + 8.13, 3.0, 3.45, 0.3, size=11,
     color=MUTED, bold=True, caps=True)
text(s, "This is the layer a competitor cannot copy by copying our prompts.\n\nIt is also "
        "what decides whether a customer can defend this work to an auditor or a "
        "regulator. Engagements die on \"show me how this entry was produced\", and that is "
        "answerable only if it was designed for at the schema level.\n\nRoughly [15–20]% of "
        "engineering effort, and almost none of the demo.",
     ML + 8.13, 3.38, 3.45, 2.7, size=12.5, color=BODY, line=1.42)

notes(s, """
The line to land: "this is the layer a competitor cannot copy by copying our prompts."

Bitemporality needs a concrete example. Use this one: a client's March books were closed in
April. In June a reviewer asks what we believed on 31 March and why. A normal database
cannot answer that — it holds current state only. We can, because we store both when a fact
became true and when we learned it. The same mechanism is what lets us take last year's real
closed books and replay our agents against them to measure accuracy, which is our entire
evaluation strategy.

If a technical investor pushes on cost: yes, append-only bitemporal ledgers are more
expensive to build and query than mutable tables. We accepted that in month one precisely
because it is unrecoverable later.

The last line of the panel is worth saying out loud — 15 to 20% of effort and almost none of
the demo. That asymmetry tells them where our engineering judgment sits.
""")


# ============================================================== 9. TOOL LAYER

s = new_slide()
eyebrow(s, "Key component · 2 of 4")
heading(s, "Tool layer: an interface for an unknown reasoner")

text(s, "This is not a REST API with a different name. A REST API assumes a client that "
        "already knows what it wants. A tool interface assumes a capable, confident, "
        "occasionally wrong caller. Five things differ:",
     ML, 1.9, 11.3, 0.65, size=14.5, color=BODY, line=1.35)

items = [
    ("Errors must be instructive, because the agent reads them and retries.", "h"),
    ("\"400 invalid\" is worthless. \"Entry unbalanced: debits 1,200.00, credits 1,150.00, "
     "difference 50.00\" produces a retry that succeeds. Error-message quality is a "
     "functional requirement here, not a nicety.", "b"),
    ("The tool owns the invariant — not the caller.", "h"),
    ("Normally you trust your own frontend to send sane data. We invert that: every tool "
     "assumes a wrong caller and enforces balance, period locks and permissions itself.", "b"),
    ("Idempotency is mandatory, because agents retry.", "h"),
    ("Granularity is a genuine tradeoff.", "h"),
    ("Too fine and the agent burns forty calls and loses the thread; too coarse and it "
     "cannot express intent. Expect to iterate — we have.", "b"),
    ("The surface must stay small and discoverable. Context is a scarce resource; you "
     "cannot expose four hundred endpoints to a reasoner.", "h"),
]
bullets(s, items, ML, 2.72, 7.4, 4.0, size=13.5, gap=7)

box(s, ML + 7.85, 2.72, 4.02, 2.3, fill=TINT_A, edge=ACCENT)
text(s, "The rule", ML + 8.13, 2.94, 3.45, 0.3, size=11, color=DEEPA, bold=True, caps=True)
text(s, "The model never does arithmetic or applies a rule in its head.\n\nEvery state "
        "change is a typed tool call that validates itself and rejects bad input. The "
        "agent's job is choosing which tool with which arguments. The tool's job is being "
        "correct.",
     ML + 8.13, 3.3, 3.45, 1.6, size=12.5, color=BODY, line=1.42)

box(s, ML + 7.85, 5.16, 4.02, 1.1, fill=PANEL)
text(s, "The single highest-leverage design decision in the system.",
     ML + 8.13, 5.42, 3.45, 0.7, size=13, color=INK, bold=True, line=1.3)

notes(s, """
The error-message point is what convinces engineers we have actually shipped this. Nobody
who has not built an agent system thinks of error strings as a functional requirement. Give
the unbalanced-entry example verbatim.

The inverted trust relationship is the second thing to emphasise, and it resonates with
anyone who understands financial controls. Every engineer has written validation in the
frontend and trusted it. You cannot do that when the caller is a model, so the invariant has
to live in the tool. Small idea, large consequences for how the codebase is organised.

The rule box is the one to repeat if you get a single point across: the model never does
arithmetic. Every wrong number in a competitor's demo comes from letting the model compute
instead of calling something that computes.

If asked about MCP or tool-protocol standardisation: we are compatible, but it is an
integration detail, not the hard part. The hard part is designing the right thirty or so
tools at the right granularity with the right invariants, and that took accounting domain
expertise, not protocol work.
""")


# ============================================================== 10. VERIFICATION

s = new_slide()
eyebrow(s, "Key component · 3 of 4")
heading(s, "Verification and confidence routing")

text(s, "The biggest change from the software we use today: from syntactic validation — is "
        "this well-formed? — to semantic verification — is this well-formed answer correct?",
     ML, 1.88, 11.4, 0.5, size=14.5, color=BODY, line=1.35)

cw3 = 3.79
xs = ML
tiers = [
    ("Tier 1 · Hard invariants", TINT_W, WARM,
     "Deterministic, non-negotiable, run before anything commits. Balance, trial-balance "
     "ties, subledger-to-control ties, reconciliation delta, duplicate detection, period "
     "locks. A failure here is a hard block, never a warning."),
    ("Tier 2 · Statistical checks", PANEL, LINE,
     "Variance against prior period, unusual vendor and account pairings, distributional "
     "anomalies. Catches the well-formed-but-wrong class that invariants pass. Produces a "
     "score, not a verdict."),
    ("Tier 3 · Routing", PANEL, LINE,
     "Above threshold: auto-post. Below: into the review queue with full context. Confidence "
     "is derived from verification signals and agreement between independent passes — never "
     "from asking the model how sure it is."),
]
for name, fill, edge, body in tiers:
    box(s, xs, 2.5, cw3, 2.7, fill=fill, edge=edge)
    text(s, name, xs + 0.26, 2.73, cw3 - 0.5, 0.32, size=13.5, color=INK, bold=True)
    text(s, body, xs + 0.26, 3.15, cw3 - 0.5, 1.95, size=12.5, color=BODY, line=1.4)
    xs += cw3 + 0.28

box(s, ML, 5.42, 6.6, 1.45, fill=PANEL2)
text(s, "Models are badly calibrated and reliably overconfident. Asking one for a confidence "
        "score is not a control.", ML + 0.3, 5.62, 6.0, 0.6, size=13.5, color=INK,
     bold=True, line=1.3)
text(s, "We derive confidence from tier-1 and tier-2 signals and from agreement between "
        "independent passes, then calibrate against known outcomes.",
     ML + 0.3, 6.24, 6.0, 0.5, size=12, color=BODY, line=1.35)

box(s, ML + 6.9, 5.42, 5.0, 1.45, fill=TINT_W)
text(s, "The threshold is the business case", ML + 7.18, 5.6, 4.45, 0.28, size=11,
     color=WARM, bold=True, caps=True)
text(s, "Move it up and hours fall but error risk rises; move it down and we have automated "
        "nothing. The entire engineering programme is raising it safely, per category, with "
        "the movement measured.",
     ML + 7.18, 5.92, 4.45, 0.85, size=12, color=DEEPW, line=1.38)

notes(s, """
This is the technical heart of the deck. If they remember one architecture slide, it should
be this one.

Open with the syntactic-to-semantic line, then the three tiers briefly. Spend your remaining
time on the bottom-left box, because it is the most common failure in the category and
stating it clearly is a strong credibility signal: models are overconfident, so
self-reported confidence is not a control. We derive confidence from verification signals
and from agreement between independent passes, and we calibrate it against known outcomes.

Then the bottom-right box, which is the commercial translation. The routing threshold IS the
business case. Every point of auto-post rate is a measurable number of hours. That is why
this is a forecast rather than a hope: we can see the rate per transaction category and we
know what each point is worth in staff time.

Likely question: "what is your current auto-post rate?" Have the real number per category
with the denominator defined. If it is early, give the number and the trend. Never give a
blended number without the denominator — a technical investor will assume you are hiding
the mix, and they will be right to.

Also worth volunteering: every reviewer correction becomes an explicit, versioned rule
first and a retrieval example second, because a rule can be shown to a reviewer and an
embedding cannot.
""")


# ============================================================== 11. EVALS

s = new_slide()
eyebrow(s, "Key component · 4 of 4")
heading(s, "Evals and replay: the measurement system is the product")

text(s, "A normal software feature is correct or it is a bug. An agent capability has a "
        "distribution. You stop asking \"does it work\" and start asking \"at what rate, on "
        "what input distribution, with what failure modes.\" That changes engineering "
        "practice more than anything else here.",
     ML, 1.9, 11.3, 0.7, size=14.5, color=BODY, line=1.35)

items = [
    ("Ground truth comes from closes that already happened.", "h"),
    ("We replay agents against historical periods and score against what the bookkeeper "
     "actually did, then against what review corrected. Our RPA exception logs tell us which "
     "categories to weight.", "b"),
    ("CI contains statistical tests.", "h"),
    ("A two-point regression may be noise or a catastrophe, and only eval volume tells you "
     "which.", "b"),
    ("Bug reports are not reproducible.", "h"),
    ("Full trace capture on every run — prompt, context, tool calls, model version — or we "
     "are blind in production.", "b"),
    ("Model upgrades are breaking changes to behaviour we never specified.", "h"),
    ("Version-pin, shadow-eval every candidate model against the corpus before any swap, "
     "keep a fallback. A standing pipeline, not a project.", "b"),
]
bullets(s, items, ML, 2.78, 7.4, 3.9, size=13.5, gap=7)

box(s, ML + 7.85, 2.78, 4.02, 3.45, fill=TINT_W, edge=WARM)
text(s, "The reporting unit changes", ML + 8.13, 3.0, 3.45, 0.3, size=11, color=WARM,
     bold=True, caps=True)
text(s, "Software roadmaps list features.\n\nOurs lists accuracy on a category:\n"
        "\"raise auto-post rate on the ambiguous-vendor bucket from 71% to 88%.\"\n\n"
        "Different planning, different staffing, and a roadmap whose progress is measured "
        "rather than asserted — which is also how we intend to report to a board.",
     ML + 8.13, 3.38, 3.45, 2.65, size=12.5, color=BODY, line=1.42)

notes(s, """
Lead with the first bullet, not the theory. The eval corpus is not a nice dataset to have;
it is the measurement instrument that makes every other claim in this deck checkable, and
it is the thing a competitor cannot clone from our marketing.

Second, be honest that evals are harder than the product. Teams underinvest here, lose the
ability to tell whether they are improving, and end up optimising on impressions. The corpus
needs a data pipeline plus senior accountant time to label — that senior time is a real cost
and it appears explicitly in the investment slide.

Our RPA history matters here too: the exception logs tell us which transaction categories
break most often and at what volume, so we weight the corpus by where the hours actually
are rather than sampling uniformly.

The panel on the right is genuinely useful to an investor because it tells them what board
reporting will look like. We will not come with "we shipped a feature." We will come with
"auto-post rate on this category moved from 71 to 88, here is what that is worth in hours."
Say that explicitly — it is a better governance relationship and it signals we know what to
measure.

The competitive point, if pressed on what stops a well-funded entrant: prompts copy in an
afternoon. A corpus of thousands of reviewed closes with the corrections attached is years
of accumulated professional work tied to client relationships. That and write permission are
the honest answers.
""")


# ============================================================== 12. HARD PROBLEMS

s = new_slide()
eyebrow(s, "Technical risk")
heading(s, "The two hard problems, stated plainly")

box(s, ML, 2.0, 5.79, 4.3, fill=PANEL, edge=LINE)
text(s, "1 · Reliability compounds badly", ML + 0.3, 2.24, 5.2, 0.35, size=17, color=INK, bold=True)
text(s, "95% per step across 20 steps is 36% end-to-end.", ML + 0.3, 2.68, 5.2, 0.55,
     size=15, color=WARM, bold=True, line=1.25)
text(s, "The central engineering problem of the category, and no prompt fixes it. The only "
        "real answer is architectural:",
     ML + 0.3, 3.26, 5.2, 0.6, size=13, color=BODY, line=1.35)
bullets(s, [("Shorten agentic spans; keep deterministic orchestration between them", "b"),
            ("Make every step independently verifiable", "b"),
            ("Gate before commit, always", "b"),
            ("Prefer many short verified hops to one long autonomous chain", "b")],
        ML + 0.3, 3.98, 5.2, 1.8, size=12.5, gap=6)
text(s, "Any plan assuming long autonomous chains without this will fail in production. Use "
        "it as a filter on anyone else you look at.",
     ML + 0.3, 5.6, 5.2, 0.6, size=12, color=WARM, bold=True, line=1.35)

box(s, ML + 6.11, 2.0, 5.79, 4.3, fill=PANEL, edge=LINE)
text(s, "2 · The security model inverts", ML + 6.41, 2.24, 5.2, 0.35, size=17, color=INK, bold=True)
text(s, "An actor with authority, taking instructions from untrusted content.",
     ML + 6.41, 2.68, 5.2, 0.55, size=15, color=WARM, bold=True, line=1.25)
text(s, "Normally code is trusted and users are untrusted. Our agent is neither — it acts "
        "with write authority while reading a vendor's PDF or an inbound email. A genuinely "
        "new vulnerability class, and a customer's books are what is exposed.",
     ML + 6.41, 3.26, 5.2, 0.9, size=13, color=BODY, line=1.35)
bullets(s, [("Scoped, least-privilege credentials per task", "b"),
            ("No single tool both reads untrusted content and takes irreversible action", "b"),
            ("Separation of duties in code: an agent cannot approve its own entry", "b"),
            ("Hard period locks; document-derived text treated as tainted", "b")],
        ML + 6.41, 4.28, 5.2, 1.8, size=12.5, gap=6)

text(s, "Two more we track and do not minimise. The 60/40 trap: the easy majority of "
        "transactions automates fast and creates false confidence while the residual holds "
        "most of the hours — so we instrument activity time, not transaction counts. And "
        "migration onto the new substrate: large, correctness-critical, and the first thing "
        "every engagement needs.",
     ML, 6.48, CW, 0.6, size=12, color=BODY, line=1.32)

notes(s, """
Do not skip this slide and do not soften it. Volunteering the hard problems with mitigations
already built is the strongest credibility move available, because every sophisticated
investor already knows about compounding error rates and is waiting to see whether we do.

On reliability: the arithmetic is the argument. 0.95^20 = 0.36. Say the number out loud.
Then make clear the fix is architectural rather than a better prompt, and that the
architecture two slides back was designed around exactly this — short spans, verification
between them, deterministic orchestration.

Offer the filter explicitly: "when you look at others in this space, ask how long their
autonomous chains run and what verifies each step. It is a fast way to tell who has shipped
something."

On security: prompt injection into an agent with ledger and bank write access is the
scenario that should worry a partner, and we should be the ones to raise it. Walk the four
mitigations. Separation of duties in code lands especially well with this audience — we took
a human control framework and applied it to a software actor.

Our own liability is the honest framing of the stakes here. It is not an abstract product
risk; it is a customer's books and our reputation. Saying so is why the controls line in the
controls budget is not negotiable.
""")


# ============================================================== 13. EFFORT SHAPE

s = new_slide()
eyebrow(s, "Engineering shape")
heading(s, "Building this is a different engineering shape")

text(s, "In conventional accounting software the system of record is the smallest, most "
        "stable, most valuable part of the codebase and the interface is the largest and most "
        "churn-prone. The ledger under a major package is a few thousand lines that barely "
        "changed in a decade; the UI around it is hundreds of thousands that changed every "
        "sprint.",
     ML, 1.88, 11.4, 0.78, size=14, color=BODY, line=1.35)

hdr = [["Layer", "Conventional", "Ours", "Nature of the change"]]
rows = [
    ["System of record", "10–15%", "15–20%", "Same concept, far harder requirements"],
    ["Tool layer (typed actions)", "—", "15–20%", "New — and not the same as a REST API"],
    ["Verification & guardrails", "—", "~15%", "New — no analogue in the software we use"],
    ["Evals & ground-truth corpus", "—", "10–15%", "New — is our test suite and our spec"],
    ["Durable orchestration", "~2%", "~10%", "Qualitatively different from cron and queues"],
    ["Agent layer (prompts, memory)", "—", "~10%", "New, and smaller than people expect"],
    ["UI", "30–40%", "15–20%", "Shrinks, and changes character entirely"],
    ["Integrations", "10–15%", "~10%", "Similar, plus agent-driven fallback"],
    ["Platform", "~15%", "~15%", "Plus agent identity and cost telemetry"],
]
table(s, hdr + rows, ML, 2.74, CW, [3.3, 1.7, 1.65, 5.24], row_h=0.35, head_h=0.42,
      size=11.5, emphasize_col=2)

box(s, ML, 6.36, CW, 0.72, fill=INK)
text(s, "Conventional software spent its engineering budget making a human effective at the "
        "interface. We spend ours making a machine's output provable at the substrate.",
     ML + 0.32, 6.52, CW - 0.64, 0.44, size=13.5, color=WHITE, bold=True, line=1.3)

notes(s, """
Percentages are directional — the shape of such codebases rather than measured data. Say so;
a technical investor will respect the caveat and stop probing the digits.

Two observations to draw out.

The interface absorbed all the heterogeneity in conventional software. Every new client
segment, edge case and regulatory variation arrived as a new screen, field or report, while
the domain model absorbed almost none of it. Interface complexity is client diversity
projected onto deterministic code. Agents absorb that diversity in the model instead, which
is why the UI share roughly halves and why we need one excellent surface rather than forty
screens.

The agent layer is one of the smallest lines. This is the third time they will hear it and
it is still the most counterintuitive thing in the deck.

The closing bar is the one-sentence summary of the technical section. Deliver it and stop —
this is the right place to take questions before the implementation plan.

If asked "is your headcount lower than a software company's?" — no, differently shaped:
fewer frontend engineers, more data and infrastructure engineering, plus our own accountants
inside the loop rather than beside it.
""")

# ============================================================== 14. HOW WE BUILD

s = new_slide()
eyebrow(s, "How we build")
heading(s, "The bottleneck is specification, not implementation")

text(s, "Coding agents changed where engineering time goes. Writing the implementation is no "
        "longer the expensive step — defining precisely what correct means, and building the "
        "harness that checks it, is. That happens to be the same asset the product needs.",
     ML, 1.88, 11.4, 0.62, size=15, color=BODY, line=1.35)

steps = [
    ("1 · Define the outcome",
     "State the invariant before the feature. \"An entry that fails trial-balance ties never "
     "commits.\" Written as an executable check, not prose in a ticket."),
    ("2 · Build the harness",
     "Fixtures from real historical closes, replay runner, diff against known outcomes, CI "
     "gate. The harness comes before the implementation it will judge."),
    ("3 · Let agents implement",
     "With a mechanically-checkable target, coding agents iterate against the harness until "
     "it passes. Volume of code stops being the constraint."),
    ("4 · Humans hold the boundaries",
     "Schema and substrate decisions, tool granularity, what counts as correct, and reviewing "
     "the gates. The judgment, not the typing."),
]
xs = ML
for name, body in steps:
    w = 2.79
    box(s, xs, 2.68, w, 2.05, fill=PANEL, edge=LINE)
    text(s, name, xs + 0.24, 2.9, w - 0.45, 0.3, size=13.5, color=ACCENT, bold=True)
    text(s, body, xs + 0.24, 3.28, w - 0.45, 1.3, size=12, color=BODY, line=1.4)
    xs += w + 0.25

box(s, ML, 4.95, 7.4, 1.95, fill=TINT_W, edge=WARM)
text(s, "Why this compresses the timeline", ML + 0.3, 5.15, 6.8, 0.28, size=10.5,
     color=WARM, bold=True, caps=True)
text(s, "The same property that makes the product work — a domain where correctness is "
        "machine-checkable — is what makes it buildable fast. A coding agent is only as good "
        "as the test it is aiming at, and accounting hands us the tests. So the plan on the "
        "next slide is measured in weeks, and it is gated on verified outcomes rather than on "
        "estimated effort.",
     ML + 0.3, 5.48, 6.8, 1.25, size=13, color=DEEPW, line=1.42)

box(s, ML + 7.65, 4.95, 4.24, 1.95, fill=PANEL2)
text(s, "The honest limit", ML + 7.93, 5.15, 3.7, 0.28, size=10.5, color=MUTED,
     bold=True, caps=True)
text(s, "Adding engineers does not compress this much. The rate limit is how fast we can "
        "specify and verify, which is domain work.\n\nThat is why the team below is small and "
        "weighted to substrate, evals and accounting expertise.",
     ML + 7.93, 5.48, 3.7, 1.25, size=12, color=BODY, line=1.4)

notes(s, """
This slide is why the timeline on the next one is credible rather than optimistic, so do not
rush it.

The argument in one line: coding agents moved the bottleneck from writing code to defining
what correct means. Anyone can now generate an implementation quickly. What they cannot do
quickly is decide what the implementation must satisfy, and prove it does. In most domains
that is the hard part and it stays hard. In ours, the domain supplies the specification —
double-entry invariants, ties, reconciliation — so we are unusually well placed to exploit
agentic coding.

Notice the recursion, and say it out loud because it is the elegant part: the harness we
build to let coding agents work IS the verification layer the product needs at runtime. The
eval corpus, the replay runner, the invariant checks — we are not building test scaffolding
and then a product. We are building the product's verification layer first and getting fast
development as a side effect.

Step 2 ordering matters: the harness comes before the implementation it will judge. If you
write the code first, you will write tests that agree with the code rather than tests that
encode what the accounting requires.

The right-hand box is the honest limit and you should volunteer it. Throwing engineers at
this does not compress the schedule, because the rate limit is specification and that is
domain work. It is also the answer to "why is your team so small" — the small team is a
consequence of the method, not a constraint we are apologising for.

If asked what could go wrong with this approach: agent-written code passing a weak test is
the failure mode. The mitigation is that our gates are accounting invariants and replay
against real closes, which are extremely hard to satisfy accidentally.
""")


# ============================================================== 15. IMPLEMENTATION PLAN

s = new_slide()
eyebrow(s, "Implementation plan")
heading(s, "Eight workstreams, gated on outcomes rather than dates")

hdr = [["Workstream", "Weeks", "The gate it must pass"]]
rows = [
    ["1 · Ledger substrate and cutover",
     "0–8",
     "Zero automated diff against the incumbent across three consecutive closes before primary flips"],
    ["2 · Migration and opening balances",
     "2–12",
     "Reconstructed trial balance ties to the prior closing TB, per entity, with no manual plug"],
    ["3 · Tool layer and invariants",
     "0–6",
     "Every state change goes through a typed tool; no path writes to the ledger directly"],
    ["4 · Verification tiers and calibration",
     "1–10",
     "Calibration error under [Y] on the top five transaction categories, measured not asserted"],
    ["5 · Eval harness and corpus",
     "0–16",
     "Every model or prompt change gated by corpus regression before it reaches a customer"],
    ["6 · Durable close orchestration",
     "4–14",
     "Full checklist resumable; no re-decision on resume, proven by replay"],
    ["7 · Exception console",
     "6–14",
     "Reviewer clears [N] items/hour with no measured loss of accuracy"],
    ["8 · Controls, audit package, SOC 2",
     "4–20",
     "One external review accepts the audit package unmodified"],
]
table(s, hdr + rows, ML, 1.95, CW, [3.5, 0.95, 7.44], row_h=0.40, head_h=0.42, size=11.5)

box(s, ML, 5.66, 5.79, 1.40, fill=TINT_W, edge=WARM)
text(s, "Critical path", ML + 0.3, 5.84, 5.2, 0.28, size=10.5, color=WARM, bold=True, caps=True)
text(s, "Workstreams 1, 3 and 5 gate everything else: the substrate, the tools that guard it, "
        "and the harness that proves it. They start in week zero and nothing downstream is "
        "meaningful until they hold.",
     ML + 0.3, 6.16, 5.2, 0.82, size=12, color=DEEPW, line=1.4)

box(s, ML + 6.11, 5.66, 5.79, 1.40, fill=PANEL)
text(s, "Roughly 20 weeks to a customer running end to end", ML + 6.41, 5.84, 5.2, 0.28,
     size=10.5, color=MUTED, bold=True, caps=True)
text(s, "Aggressive by convention, and it rests on the previous slide: the gates are "
        "machine-checkable, so implementation iterates against them without waiting on us. "
        "Slippage shows up as a failed gate, not a status report.",
     ML + 6.41, 6.16, 5.2, 0.82, size=12, color=BODY, line=1.4)

notes(s, """
The gate column is the point of this slide. Every workstream has a testable exit condition
rather than a delivery date, which is how you should ask us to govern it. Read two or three
gates out loud — they are unusually concrete for a technology plan and that is the impression
to leave.

On workstream 1, the detail that matters is the dual-write diff harness. We do not flip the
primary ledger on a judgment call. We run both systems in parallel, diff every entry
automatically, and require zero diff across three consecutive closes before switching. A
testable gate, not a decision.

On workstream 2, say plainly that migration is the most underestimated item in this kind of
programme. Reconstructing an opening trial balance that ties to the prior closing balance,
per entity, is unglamorous and correctness-critical, and every customer needs it before
anything else works. Imported records carry a distinct provenance class so we never claim
verified lineage we do not have.

Workstream 6's checkpointing detail is worth ten seconds with a technical audience: on
resume, a durable workflow must not silently re-run a model call and get a different answer.
We checkpoint the decision, not just the step. Very few teams get this right.

On the 20 weeks: expect scepticism and welcome it. The answer is not "we work hard." It is
that our gates are executable, so implementation iterates against a machine rather than
against review cycles. Then add the honest part — the weeks that can slip are 2 and 8,
because migration depends on the state of customer data we do not control, and SOC 2 depends
on an external auditor's calendar.

Workstream 8 exists because customer trust is the asset at risk. Do not let it be read as
overhead.
""")


# ============================================================== 16. RESOURCING

s = new_slide()
eyebrow(s, "Resourcing")
heading(s, "A small team, weighted away from implementation")

text(s, "Illustrative shape rather than a hiring plan. The notable feature is what is small: "
        "the agent layer, and implementation headcount generally.",
     ML, 1.9, 11.4, 0.4, size=14.5, color=BODY, line=1.35)

hdr = [["Function", "FTE", "Why it is sized this way"]]
rows = [
    ["Ledger substrate and tool layer", "2",
     "Schema, bitemporality, provenance, typed invariant-enforcing actions. Irreversible decisions, so senior."],
    ["Evals, harness and data", "2",
     "Corpus, replay runner, CI gates, ingestion quality, entity resolution. The largest engineering line, deliberately."],
    ["Verification and calibration", "1.5",
     "Invariants, anomaly detection, confidence derived from signals rather than self-report."],
    ["Orchestration and exception console", "1.5",
     "Durable workflows and one surface built well, instead of forty screens."],
    ["Accounting domain expertise", "2–3",
     "Writes the specifications and labels the corpus. This is the rate limit on the whole plan."],
    ["Security and controls", "1",
     "Agent identity, least privilege, separation of duties, audit package, SOC 2."],
    ["Agent layer", "1",
     "Prompts, decomposition, policy memory. The cheapest and most replaceable component."],
]
table(s, hdr + rows, ML, 2.5, CW, [3.5, 0.75, 7.64], row_h=0.42, head_h=0.42, size=11.5)

box(s, ML, 5.98, 3.79, 1.05, fill=PANEL)
text(s, "~11 FTE", ML + 0.26, 6.14, 3.3, 0.32, size=15, color=ACCENT, bold=True)
text(s, "Two of those eleven touch prompts or implementation volume.",
     ML + 0.26, 6.5, 3.3, 0.44, size=11.5, color=BODY, line=1.32)

box(s, ML + 4.07, 5.98, 3.79, 1.05, fill=PANEL)
text(s, "Non-headcount", ML + 4.33, 6.14, 3.3, 0.28, size=10.5, color=MUTED, bold=True, caps=True)
text(s, "Inference is COGS, not R&D · SOC 2 Type II · E&O review · data infra · "
        "customer consent updates",
     ML + 4.33, 6.44, 3.3, 0.5, size=10.5, color=BODY, line=1.3)

box(s, ML + 8.14, 5.98, 3.75, 1.05, fill=TINT_W, edge=WARM)
text(s, "To fill in", ML + 8.4, 6.14, 3.25, 0.28, size=10.5, color=WARM, bold=True, caps=True)
text(s, "[Total investment over 20 weeks] and [the labour cost it removes]. "
        "Do not present brackets live.",
     ML + 8.4, 6.44, 3.25, 0.5, size=10.5, color=DEEPW, line=1.3)

notes(s, """
Three things to draw out.

The agent layer is 1 of ~11 FTE. Point at it. It is the fourth time in the deck that the
demo-able part turns out to be the small part, and by now it should read as discipline rather
than as an oversight. Evals and harness is twice the size of the agent layer, and that ratio
is the single best summary of how we think.

Accounting domain expertise is the rate limit, not engineering capacity. Say it directly:
if you gave us double the engineers tomorrow, the plan would not finish meaningfully sooner,
because the constraint is how fast we can specify what correct means and label the corpus.
That is also why we are not asking for a large team.

Senior accountant time for labelling is a real, easily-forgotten cost. Budget it explicitly
rather than pretending it is free, because if you do not, the corpus quietly does not get
built and every metric in this deck becomes unverifiable.

Inference sits in COGS, not R&D. Most plans in this category get this wrong, and getting it
right is what makes the cost-per-close comparison honest.

Expected question: "is eleven people really enough for a ledger, an agent system and SOC 2 in
twenty weeks?" The honest answer has two halves. Yes for implementation volume, because
coding agents cover that against executable gates. And the thing that would break first is
not code — it is specification throughput, which is why domain expertise is the line we would
grow first if we are behind.

[FILL BEFORE USE] total investment and the loaded-cost comparison.
""")


# ============================================================== 17. METRICS

s = new_slide()
eyebrow(s, "How we know it worked")
heading(s, "The metrics, and the definition of done")

xs = ML
mets = [
    ("Auto-post rate", "Share of transactions posted with zero human touch, reported by "
     "category with the denominator stated. Product quality and hours saved in one number."),
    ("Hours per close,\nper entity", "The honest automation measure. Transaction counts "
     "flatter us; hours do not, because the residual tail holds most of the cost."),
    ("Cost per close vs.\nlabour replaced", "Inference plus review, against the fully-loaded "
     "cost of the work it removes. If this does not fall, nothing else matters."),
    ("Restatement /\nerror rate", "The metric that can end this. Governed against an "
     "absolute ceiling, not a trend — and measured against the current baseline."),
]
for name, body in mets:
    w = 2.79
    box(s, xs, 1.9, w, 2.3, fill=PANEL)
    text(s, name, xs + 0.24, 2.11, w - 0.45, 0.7, size=15.5, color=ACCENT, bold=True, line=1.15)
    text(s, body, xs + 0.24, 2.92, w - 0.45, 1.2, size=12, color=BODY, line=1.4)
    xs += w + 0.25

text(s, "The baseline is not zero error. Today's RPA-plus-human process makes mistakes, and "
        "we can measure its rate. That is the bar an agent has to beat — not perfection.",
     ML, 4.34, CW, 0.4, size=13.5, color=WARM, bold=True, line=1.3)

box(s, ML, 4.8, CW, 1.58, fill=PANEL2)
text(s, "Definition of done", ML + 0.32, 4.98, 5.0, 0.28, size=10.5, color=MUTED,
     bold=True, caps=True)
gates_l = [("Our ledger primary for [X] customers, zero dual-write diff for three closes", "b"),
           ("Auto-post rate at or above [X]% on the top five categories, denominators stated", "b"),
           ("Hours per close per entity down [Z]% against the current baseline", "b")]
gates_r = [("Restatement rate at or below the pre-agent baseline", "b"),
           ("Cost per close below [W]% of the labour cost it replaces", "b"),
           ("SOC 2 Type II achieved; audit package accepted in one external review", "b")]
bullets(s, gates_l, ML + 0.32, 5.32, 5.6, 1.0, size=12, gap=5)
bullets(s, gates_r, ML + 6.2, 5.32, 5.4, 1.0, size=12, gap=5)

box(s, ML, 6.48, CW, 0.6, fill=INK)
text(s, "Not AI features on top of accounting software. Software that does the accounting — "
        "on a substrate where every number can be proven and defended.",
     ML + 0.32, 6.6, CW - 0.64, 0.42, size=13, color=WHITE, bold=True, line=1.3)

notes(s, """
Close by handing them the scorecard. Offering the metrics you can be judged on — including
the one that can end the programme — is a stronger close than a projection, and it sets up
the governance relationship we actually want.

On auto-post rate: always with the denominator and the category mix. A blended number with no
denominator is the standard way companies in this category flatter themselves, and a
technical investor will assume that is what is happening unless you pre-empt it.

The line above the definition-of-done box is the one to say slowly, and our RPA history is
what lets us say it credibly. The incumbent process is not error-free. Bots misfire, humans
miscode, and closes get corrected. We have the logs. So the honest bar for an agent is the
measured error rate of what customers do today, and we intend to publish both numbers. That
converts an unbounded fear into a comparison — and very few teams can produce the baseline.

On restatement rate: note deliberately that this one is governed against an absolute ceiling,
not a trend. Everything else we optimise; this one we bound. That distinction is what a
finance leader needs to hear before granting write access to a ledger.

Final line is the positioning statement. Say it, stop, take questions.

Backup slides worth having ready: cost-per-close model at current token prices, gross-margin
slope as automation rate rises, the RPA telemetry on intervention rates by category,
competitive landscape, and the liability and consent structure.
""")


prs.save("/Users/cllu/Projects/aimi/agent-native-accounting-deck.pptx")
print(f"saved · {_page['n'] + 1} slides")
